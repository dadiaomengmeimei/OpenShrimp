"""
Shared File Toolkit – common file operations for all sub-apps.

Provides:
- PDF parsing (extract text from PDF)
- PDF generation (create PDF from text/HTML)
- PPT generation (create PowerPoint presentations)
- Excel parsing (read Excel/CSV into structured data)
- Excel generation (create Excel from data)
- File download URL management (register files and get download tokens)

Usage in sub-apps:
    from backend.core.file_toolkit import (
        parse_pdf, generate_pdf,
        generate_ppt,
        parse_excel, generate_excel,
        register_download, get_download_url,
    )
"""
from __future__ import annotations

import csv
import io
import json
import os
import secrets
import tempfile
import time
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

# ---------------------------------------------------------------------------
# Download registry – maps tokens to file paths for secure download
# ---------------------------------------------------------------------------

# {token: {"path": str, "filename": str, "mime": str, "created_at": float, "ttl": int}}
_download_registry: Dict[str, Dict[str, Any]] = {}

# Default time-to-live for download tokens (seconds)
DEFAULT_TTL = 3600 * 24  # 24 hours

# Shared output directory for generated files
_OUTPUT_DIR = Path(__file__).resolve().parent.parent / "generated_files"
_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def register_download(
    file_path: Union[str, Path],
    *,
    filename: Optional[str] = None,
    mime_type: Optional[str] = None,
    ttl: int = DEFAULT_TTL,
) -> str:
    """
    Register a file for download and return a unique token.

    Args:
        file_path: Absolute or relative path to the generated file.
        filename: Display filename for download (defaults to original name).
        mime_type: MIME type (auto-detected if not provided).
        ttl: Time-to-live in seconds for the download token.

    Returns:
        A unique token string that can be used to construct the download URL.
    """
    path = Path(file_path).resolve()
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")

    if not filename:
        filename = path.name

    if not mime_type:
        mime_type = _guess_mime(path.suffix)

    token = secrets.token_urlsafe(16)
    _download_registry[token] = {
        "path": str(path),
        "filename": filename,
        "mime": mime_type,
        "created_at": time.time(),
        "ttl": ttl,
    }

    # Clean up expired tokens periodically
    _cleanup_expired()

    return token


def get_download_url(token: str, base_url: str = "") -> str:
    """
    Build a full download URL from a token.

    Args:
        token: The download token from register_download().
        base_url: Optional base URL (e.g. "http://localhost:8000").
                  If empty, returns just the relative path.

    Returns:
        The full or relative download URL.
    """
    relative = f"/api/files/download/{token}"
    if base_url:
        return f"{base_url.rstrip('/')}{relative}"
    return relative


def get_download_info(token: str) -> Optional[Dict[str, Any]]:
    """
    Look up a download token and return file info if valid.

    Returns None if the token is expired or not found.
    """
    info = _download_registry.get(token)
    if not info:
        return None
    if time.time() - info["created_at"] > info["ttl"]:
        del _download_registry[token]
        return None
    return info


def _cleanup_expired():
    """Remove expired download tokens."""
    now = time.time()
    expired = [k for k, v in _download_registry.items() if now - v["created_at"] > v["ttl"]]
    for k in expired:
        del _download_registry[k]


def _guess_mime(suffix: str) -> str:
    """Guess MIME type from file extension."""
    mime_map = {
        ".pdf": "application/pdf",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".csv": "text/csv",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".txt": "text/plain",
        ".json": "application/json",
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".zip": "application/zip",
    }
    return mime_map.get(suffix.lower(), "application/octet-stream")


# ---------------------------------------------------------------------------
# PDF Operations
# ---------------------------------------------------------------------------

def parse_pdf(file_path: Union[str, Path, io.BytesIO]) -> str:
    """
    Extract text content from a PDF file.

    Args:
        file_path: Path to PDF file, or a BytesIO object.

    Returns:
        Extracted text as a single string.
    """
    from PyPDF2 import PdfReader

    if isinstance(file_path, io.BytesIO):
        reader = PdfReader(file_path)
    else:
        reader = PdfReader(str(file_path))

    pages_text = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages_text.append(f"--- Page {i + 1} ---\n{text.strip()}")

    return "\n\n".join(pages_text)


def generate_pdf(
    content: str,
    output_path: Optional[Union[str, Path]] = None,
    *,
    title: Optional[str] = None,
) -> Path:
    """
    Generate a simple PDF from text content.

    Uses reportlab if available, otherwise falls back to a basic approach.

    Args:
        content: Text content to put in the PDF.
        output_path: Where to save the PDF. Auto-generated if not provided.
        title: Optional document title.

    Returns:
        Path to the generated PDF file.
    """
    if output_path is None:
        output_path = _OUTPUT_DIR / f"doc_{secrets.token_hex(6)}.pdf"
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import cm
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        doc = SimpleDocTemplate(str(output_path), pagesize=A4,
                                leftMargin=2 * cm, rightMargin=2 * cm,
                                topMargin=2 * cm, bottomMargin=2 * cm)
        styles = getSampleStyleSheet()

        # Try to register a CJK font for Chinese support
        _try_register_cjk_font()

        body_style = ParagraphStyle(
            'Body', parent=styles['Normal'],
            fontSize=11, leading=16,
            fontName=_get_cjk_font_name() or 'Helvetica',
        )
        title_style = ParagraphStyle(
            'DocTitle', parent=styles['Title'],
            fontSize=18, leading=24,
            fontName=_get_cjk_font_name() or 'Helvetica-Bold',
        )

        story = []
        if title:
            story.append(Paragraph(title, title_style))
            story.append(Spacer(1, 0.5 * cm))

        for line in content.split('\n'):
            if line.strip():
                # Escape XML special chars for reportlab
                safe_line = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                story.append(Paragraph(safe_line, body_style))
            else:
                story.append(Spacer(1, 0.3 * cm))

        doc.build(story)
        print(f"[file_toolkit] PDF generated: {output_path}")

    except ImportError:
        # Fallback: use FPDF2 if available, else raise helpful error
        try:
            from fpdf import FPDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Helvetica", size=11)
            if title:
                pdf.set_font("Helvetica", "B", 16)
                pdf.cell(0, 10, title, ln=True, align="C")
                pdf.ln(5)
                pdf.set_font("Helvetica", size=11)
            pdf.multi_cell(0, 7, content)
            pdf.output(str(output_path))
            print(f"[file_toolkit] PDF generated (fpdf2 fallback): {output_path}")
        except ImportError:
            raise ImportError(
                "PDF generation requires 'reportlab' or 'fpdf2'. "
                "Add one of them to your app's requirements.txt."
            )

    return output_path


def _try_register_cjk_font():
    """Try to register a CJK-compatible font for Chinese/Japanese/Korean text."""
    try:
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        # Common CJK font paths on different OS
        font_paths = [
            # macOS
            "/System/Library/Fonts/STHeiti Light.ttc",
            "/System/Library/Fonts/PingFang.ttc",
            "/Library/Fonts/Arial Unicode.ttf",
            # Linux
            "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
            "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
            # Windows
            "C:/Windows/Fonts/msyh.ttc",
            "C:/Windows/Fonts/simsun.ttc",
        ]
        for fp in font_paths:
            if os.path.exists(fp):
                try:
                    pdfmetrics.registerFont(TTFont("CJKFont", fp))
                    return
                except Exception:
                    continue
    except Exception:
        pass


def _get_cjk_font_name() -> Optional[str]:
    """Return CJK font name if registered, else None."""
    try:
        from reportlab.pdfbase import pdfmetrics
        if "CJKFont" in pdfmetrics.getRegisteredFontNames():
            return "CJKFont"
    except Exception:
        pass
    return None


# ---------------------------------------------------------------------------
# PPT Operations
# ---------------------------------------------------------------------------

def generate_ppt(
    slides: List[Dict[str, Any]],
    output_path: Optional[Union[str, Path]] = None,
    *,
    title: Optional[str] = None,
    style: str = "professional",
) -> Path:
    """
    Generate a PowerPoint presentation from structured slide data.

    Args:
        slides: List of dicts, each with:
            - "title": str (slide title)
            - "content": list[str] (bullet points)
            - "notes": str (optional speaker notes)
        output_path: Where to save the PPTX. Auto-generated if not provided.
        title: Overall presentation title (used for the first slide).
        style: Visual style ("professional", "creative", "minimal", "academic").

    Returns:
        Path to the generated PPTX file.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    if output_path is None:
        output_path = _OUTPUT_DIR / f"ppt_{secrets.token_hex(6)}.pptx"
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    colors = _get_ppt_style(style)

    # Title slide
    if title:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        _add_bg(slide, prs, colors["bg"])
        _add_accent_bar(slide, prs, colors["accent"], top=True)

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*colors["title"])
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Content slides
    for idx, s in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_bg(slide, prs, colors["bg"])
        _add_accent_bar(slide, prs, colors["accent"], top=True)

        # Slide title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.733), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = s.get("title", f"Slide {idx + 1}")
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*colors["title"])

        # Bullet points
        content_items = s.get("content", [])
        if content_items:
            txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.733), Inches(4.5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for point in content_items:
                p2 = tf2.add_paragraph()
                p2.text = f"• {point}"
                p2.font.size = Pt(22)
                p2.font.color.rgb = RGBColor(*colors["content"])
                p2.space_before = Pt(8)
                p2.line_spacing = 1.3

        # Speaker notes
        notes_text = s.get("notes", "")
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    prs.save(str(output_path))
    print(f"[file_toolkit] PPT generated: {output_path} ({len(slides)} slides)")
    return output_path


def _get_ppt_style(style: str) -> Dict[str, tuple]:
    styles = {
        "professional": {"title": (25, 25, 112), "content": (50, 50, 50), "bg": (255, 255, 255), "accent": (41, 128, 185)},
        "creative": {"title": (220, 20, 60), "content": (60, 60, 60), "bg": (255, 250, 240), "accent": (230, 126, 34)},
        "minimal": {"title": (0, 0, 0), "content": (64, 64, 64), "bg": (255, 255, 255), "accent": (0, 0, 0)},
        "academic": {"title": (0, 100, 0), "content": (40, 40, 40), "bg": (255, 255, 250), "accent": (39, 174, 96)},
    }
    return styles.get(style, styles["professional"])


def _add_bg(slide, prs, color: tuple):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(*color)
    bg.line.fill.background()


def _add_accent_bar(slide, prs, color: tuple, *, top: bool = True):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    y = Inches(0) if top else Inches(7.0)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), y, prs.slide_width, Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*color)
    bar.line.fill.background()


# ---------------------------------------------------------------------------
# Excel / CSV Operations
# ---------------------------------------------------------------------------

def parse_excel(
    file_path: Union[str, Path, io.BytesIO],
    *,
    sheet_name: Optional[str] = None,
    max_rows: int = 10000,
) -> Dict[str, Any]:
    """
    Parse an Excel or CSV file into structured data.

    Args:
        file_path: Path to .xlsx/.csv file, or BytesIO.
        sheet_name: Specific sheet to read (Excel only). None = first sheet.
        max_rows: Maximum rows to read.

    Returns:
        Dict with keys: "headers" (list[str]), "rows" (list[list]), "row_count" (int),
        "sheet_names" (list[str], Excel only).
    """
    import pandas as pd

    path = file_path if isinstance(file_path, io.BytesIO) else Path(file_path)
    suffix = "" if isinstance(path, io.BytesIO) else path.suffix.lower()

    if suffix == ".csv" or (isinstance(path, io.BytesIO) and not sheet_name):
        # Try CSV first
        try:
            df = pd.read_csv(path, nrows=max_rows)
            return {
                "headers": list(df.columns),
                "rows": df.values.tolist(),
                "row_count": len(df),
                "sheet_names": [],
            }
        except Exception:
            pass

    # Excel
    xls = pd.ExcelFile(path)
    target_sheet = sheet_name or xls.sheet_names[0]
    df = pd.read_excel(xls, sheet_name=target_sheet, nrows=max_rows)

    return {
        "headers": list(df.columns),
        "rows": df.values.tolist(),
        "row_count": len(df),
        "sheet_names": xls.sheet_names,
    }


def generate_excel(
    data: List[Dict[str, Any]],
    output_path: Optional[Union[str, Path]] = None,
    *,
    sheet_name: str = "Sheet1",
    headers: Optional[List[str]] = None,
) -> Path:
    """
    Generate an Excel file from structured data.

    Args:
        data: List of dicts (each dict = one row) or list of lists.
        output_path: Where to save. Auto-generated if not provided.
        sheet_name: Name of the worksheet.
        headers: Column headers (auto-detected from dict keys if not provided).

    Returns:
        Path to the generated .xlsx file.
    """
    import pandas as pd

    if output_path is None:
        output_path = _OUTPUT_DIR / f"data_{secrets.token_hex(6)}.xlsx"
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    if data and isinstance(data[0], dict):
        df = pd.DataFrame(data)
        if headers:
            df = df[headers]  # Reorder columns
    elif data and isinstance(data[0], (list, tuple)):
        df = pd.DataFrame(data, columns=headers)
    else:
        df = pd.DataFrame()

    df.to_excel(str(output_path), sheet_name=sheet_name, index=False)
    print(f"[file_toolkit] Excel generated: {output_path} ({len(df)} rows)")
    return output_path


def generate_csv(
    data: List[Dict[str, Any]],
    output_path: Optional[Union[str, Path]] = None,
    *,
    headers: Optional[List[str]] = None,
) -> Path:
    """
    Generate a CSV file from structured data.

    Args:
        data: List of dicts or list of lists.
        output_path: Where to save. Auto-generated if not provided.
        headers: Column headers.

    Returns:
        Path to the generated .csv file.
    """
    if output_path is None:
        output_path = _OUTPUT_DIR / f"data_{secrets.token_hex(6)}.csv"
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    with open(output_path, "w", newline="", encoding="utf-8-sig") as f:
        if data and isinstance(data[0], dict):
            keys = headers or list(data[0].keys())
            writer = csv.DictWriter(f, fieldnames=keys)
            writer.writeheader()
            writer.writerows(data)
        elif data and isinstance(data[0], (list, tuple)):
            writer = csv.writer(f)
            if headers:
                writer.writerow(headers)
            writer.writerows(data)

    print(f"[file_toolkit] CSV generated: {output_path}")
    return output_path


# ---------------------------------------------------------------------------
# Word Document Operations
# ---------------------------------------------------------------------------

def parse_docx(file_path: Union[str, Path, io.BytesIO]) -> str:
    """
    Extract text content from a .docx Word document.

    Args:
        file_path: Path to .docx file, or a BytesIO object.

    Returns:
        Extracted text as a string.
    """
    from docx import Document

    doc = Document(file_path if isinstance(file_path, io.BytesIO) else str(file_path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


# ---------------------------------------------------------------------------
# Download Link Helpers (THE standard way to serve files to users)
# ---------------------------------------------------------------------------

def make_download_link(
    file_path: Union[str, Path],
    *,
    label: Optional[str] = None,
    filename: Optional[str] = None,
) -> str:
    """
    One-step helper: register a file for download and return a **Markdown link**
    that the frontend will render as a clickable download button.

    This is THE recommended way to provide file downloads in any app.
    It returns a relative-path markdown link like:
        [📥 Download report.pdf](/api/files/download/abc123...)

    The frontend's ReactMarkdown renderer recognises URLs matching
    `/api/files/download/` and renders them as styled download buttons.

    Args:
        file_path: Path to the file on disk.
        label:     Display text for the link (default: "下载 <filename>").
        filename:  Filename shown to user when downloading (default: original name).

    Returns:
        A Markdown link string ready to embed in the reply content.

    Example::

        from backend.core.file_toolkit import make_download_link

        link = make_download_link("/tmp/report.pdf", label="下载报告")
        return {"content": f"文件已生成！\n\n{link}"}
    """
    path = Path(file_path).resolve()
    display = filename or path.name
    token = register_download(path, filename=display)
    url = get_download_url(token)  # relative: /api/files/download/{token}
    link_label = label or f"下载 {display}"
    return f"[📥 {link_label}]({url})"


def get_preview_url(token: str, base_url: str = "") -> str:
    """
    Build a preview (inline) URL from a token.
    Unlike get_download_url(), the preview endpoint serves files with
    Content-Disposition: inline, so browsers display them (e.g. images)
    instead of triggering a download.

    Args:
        token: The download token from register_download().
        base_url: Optional base URL. If empty, returns relative path.

    Returns:
        The full or relative preview URL.
    """
    relative = f"/api/files/preview/{token}"
    if base_url:
        return f"{base_url.rstrip('/')}{relative}"
    return relative


def make_preview_link(
    file_path: Union[str, Path],
    *,
    label: Optional[str] = None,
    filename: Optional[str] = None,
) -> str:
    """
    Register a file and return a Markdown link to the **preview** endpoint.
    The preview endpoint serves files inline (Content-Disposition: inline)
    instead of forcing a download. Useful for PDFs, text files, etc.

    Args:
        file_path: Path to the file on disk.
        label:     Display text for the link.
        filename:  Display filename.

    Returns:
        A Markdown link string like: [🔍 Preview report.pdf](/api/files/preview/xxx)
    """
    path = Path(file_path).resolve()
    display = filename or path.name
    token = register_download(path, filename=display)
    url = get_preview_url(token)
    link_label = label or f"预览 {display}"
    return f"[🔍 {link_label}]({url})"


def make_image_embed(
    file_path: Union[str, Path],
    *,
    alt_text: Optional[str] = None,
    filename: Optional[str] = None,
) -> str:
    """
    Register an image file and return a **Markdown image** tag that the
    frontend will render inline in the chat message.

    Uses the preview endpoint (Content-Disposition: inline) so that the
    browser fetches and displays the image directly.

    This is THE recommended way to show charts/images inline in chat.

    Args:
        file_path: Path to the image file (PNG, JPG, etc.).
        alt_text:  Alt text for the image (default: filename).
        filename:  Display filename for the image.

    Returns:
        A Markdown image string like: ![chart](/api/files/preview/xxx)

    Example::

        from backend.core.file_toolkit import generate_chart, make_image_embed

        path = generate_chart("bar", data, title="Sales")
        img = make_image_embed(path, alt_text="Sales Chart")
        return {"content": f"Here is the chart:\n\n{img}"}
    """
    path = Path(file_path).resolve()
    display = filename or path.name
    token = register_download(path, filename=display)
    url = get_preview_url(token)
    alt = alt_text or display
    return f"![{alt}]({url})"


# ---------------------------------------------------------------------------
# Convenience: generate + register in one step
# ---------------------------------------------------------------------------

def generate_and_register_pdf(
    content: str,
    *,
    title: Optional[str] = None,
    filename: Optional[str] = None,
) -> Dict[str, str]:
    """Generate a PDF and register it for download.
    
    Returns dict with keys: token, url, path, markdown_link.
    The `markdown_link` value is a ready-to-use Markdown download link.
    """
    path = generate_pdf(content, title=title)
    display = filename or path.name
    token = register_download(path, filename=display)
    url = get_download_url(token)
    return {"token": token, "url": url, "path": str(path), "markdown_link": f"[📥 下载 {display}]({url})"}


def generate_and_register_ppt(
    slides: List[Dict[str, Any]],
    *,
    title: Optional[str] = None,
    style: str = "professional",
    filename: Optional[str] = None,
) -> Dict[str, str]:
    """Generate a PPT and register it for download.
    
    Returns dict with keys: token, url, path, markdown_link.
    The `markdown_link` value is a ready-to-use Markdown download link.
    """
    path = generate_ppt(slides, title=title, style=style)
    display = filename or path.name
    token = register_download(path, filename=display)
    url = get_download_url(token)
    return {"token": token, "url": url, "path": str(path), "markdown_link": f"[📥 下载 {display}]({url})"}


def generate_and_register_excel(
    data: List[Dict[str, Any]],
    *,
    sheet_name: str = "Sheet1",
    headers: Optional[List[str]] = None,
    filename: Optional[str] = None,
) -> Dict[str, str]:
    """Generate an Excel and register it for download.
    
    Returns dict with keys: token, url, path, markdown_link.
    The `markdown_link` value is a ready-to-use Markdown download link.
    """
    path = generate_excel(data, sheet_name=sheet_name, headers=headers)
    display = filename or path.name
    token = register_download(path, filename=display)
    url = get_download_url(token)
    return {"token": token, "url": url, "path": str(path), "markdown_link": f"[📥 下载 {display}]({url})"}


def register_existing_file(
    file_path: Union[str, Path],
    *,
    filename: Optional[str] = None,
) -> Dict[str, str]:
    """Register any existing file for download.
    
    Returns dict with keys: token, url, path, markdown_link.
    The `markdown_link` value is a ready-to-use Markdown download link.
    """
    path = Path(file_path).resolve()
    display = filename or path.name
    token = register_download(path, filename=display)
    url = get_download_url(token)
    return {"token": token, "url": url, "path": str(path), "markdown_link": f"[📥 下载 {display}]({url})"}


# ---------------------------------------------------------------------------
# Text / String Utilities
# ---------------------------------------------------------------------------

def truncate_text(text: str, max_length: int = 500, suffix: str = "...") -> str:
    """
    Truncate text to a max length, appending a suffix if truncated.
    Cuts at word boundaries when possible.
    """
    if len(text) <= max_length:
        return text
    cut = max_length - len(suffix)
    # Try to cut at a space boundary
    space_idx = text.rfind(" ", 0, cut)
    if space_idx > cut * 0.6:
        cut = space_idx
    return text[:cut] + suffix


def extract_json_from_text(text: str) -> Optional[Any]:
    """
    Extract the first valid JSON object or array from a text string.
    Useful for parsing LLM responses that embed JSON in markdown or prose.

    Returns the parsed JSON value, or None if no valid JSON found.
    """
    import re

    # Strip markdown code blocks
    cleaned = re.sub(r"```(?:json)?\s*", "", text)
    cleaned = re.sub(r"```\s*", "", cleaned)

    # Try to find a JSON array
    arr_match = re.search(r"\[.*\]", cleaned, re.DOTALL)
    if arr_match:
        try:
            return json.loads(arr_match.group())
        except json.JSONDecodeError:
            pass

    # Try to find a JSON object
    obj_match = re.search(r"\{.*\}", cleaned, re.DOTALL)
    if obj_match:
        try:
            return json.loads(obj_match.group())
        except json.JSONDecodeError:
            pass

    # Try the whole thing
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        return None


def sanitize_filename(name: str, max_length: int = 100) -> str:
    """
    Sanitize a string to be safe as a filename.
    Removes or replaces characters that are not allowed in filenames.
    """
    import re
    # Replace path separators and special chars
    safe = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "_", name)
    # Collapse multiple underscores
    safe = re.sub(r"_+", "_", safe).strip("_. ")
    return safe[:max_length] if safe else "untitled"


# ---------------------------------------------------------------------------
# Image Utilities
# ---------------------------------------------------------------------------

def generate_chart(
    chart_type: str,
    data: Dict[str, Any],
    output_path: Optional[Union[str, Path]] = None,
    *,
    title: Optional[str] = None,
    figsize: tuple = (10, 6),
) -> Path:
    """
    Generate a chart image using matplotlib.

    Args:
        chart_type: One of "bar", "line", "pie", "scatter", "histogram".
        data: Dict with chart-specific keys:
            - bar/line: {"labels": [...], "values": [...], "ylabel": "...", "xlabel": "..."}
            - pie:      {"labels": [...], "values": [...]}
            - scatter:  {"x": [...], "y": [...], "xlabel": "...", "ylabel": "..."}
            - histogram:{"values": [...], "bins": 20, "xlabel": "..."}
        output_path: Where to save the image. Auto-generated if not provided.
        title: Optional chart title.
        figsize: Figure size as (width, height) in inches.

    Returns:
        Path to the generated .png file.
    """
    import matplotlib
    matplotlib.use("Agg")  # Non-interactive backend
    import matplotlib.pyplot as plt

    if output_path is None:
        output_path = _OUTPUT_DIR / f"chart_{secrets.token_hex(6)}.png"
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Enable CJK font support
    plt.rcParams["font.sans-serif"] = [
        "PingFang SC", "Heiti SC", "Microsoft YaHei",
        "WenQuanYi Zen Hei", "DejaVu Sans",
    ]
    plt.rcParams["axes.unicode_minus"] = False

    fig, ax = plt.subplots(figsize=figsize)

    if chart_type == "bar":
        ax.bar(data["labels"], data["values"], color="#4A90D9")
        ax.set_xlabel(data.get("xlabel", ""))
        ax.set_ylabel(data.get("ylabel", ""))
        plt.xticks(rotation=45, ha="right")

    elif chart_type == "line":
        ax.plot(data["labels"], data["values"], marker="o", color="#4A90D9", linewidth=2)
        ax.set_xlabel(data.get("xlabel", ""))
        ax.set_ylabel(data.get("ylabel", ""))
        plt.xticks(rotation=45, ha="right")

    elif chart_type == "pie":
        ax.pie(data["values"], labels=data["labels"], autopct="%1.1f%%", startangle=90)
        ax.axis("equal")

    elif chart_type == "scatter":
        ax.scatter(data["x"], data["y"], alpha=0.7, color="#4A90D9")
        ax.set_xlabel(data.get("xlabel", ""))
        ax.set_ylabel(data.get("ylabel", ""))

    elif chart_type == "histogram":
        ax.hist(data["values"], bins=data.get("bins", 20), color="#4A90D9", edgecolor="white")
        ax.set_xlabel(data.get("xlabel", ""))
        ax.set_ylabel(data.get("ylabel", "Frequency"))

    else:
        raise ValueError(f"Unsupported chart type: {chart_type}")

    if title:
        ax.set_title(title, fontsize=14, fontweight="bold")

    plt.tight_layout()
    fig.savefig(str(output_path), dpi=150, bbox_inches="tight")
    plt.close(fig)

    print(f"[file_toolkit] Chart generated: {output_path}")
    return output_path


def generate_and_register_chart(
    chart_type: str,
    data: Dict[str, Any],
    *,
    title: Optional[str] = None,
    filename: Optional[str] = None,
) -> Dict[str, str]:
    """Generate a chart image and register for download.
    
    Returns dict with keys: token, url, path, markdown_link, preview_url, image_embed.
    - `markdown_link`: a Markdown download link (forces download).
    - `image_embed`: a Markdown image tag that renders the chart inline in chat.
    - `preview_url`: the inline-preview URL.
    """
    path = generate_chart(chart_type, data, title=title)
    display = filename or path.name
    token = register_download(path, filename=display)
    url = get_download_url(token)
    preview = get_preview_url(token)
    alt = title or display
    return {
        "token": token,
        "url": url,
        "path": str(path),
        "markdown_link": f"[📥 下载 {display}]({url})",
        "preview_url": preview,
        "image_embed": f"![{alt}]({preview})",
    }


# ---------------------------------------------------------------------------
# Markdown / HTML Utilities
# ---------------------------------------------------------------------------

def markdown_to_html(md_text: str) -> str:
    """
    Convert Markdown text to HTML.
    Useful for rendering rich content in PDFs or emails.

    Falls back to basic conversion if markdown library is not available.
    """
    try:
        import markdown
        return markdown.markdown(md_text, extensions=["tables", "fenced_code", "nl2br"])
    except ImportError:
        import re
        # Basic fallback: just handle paragraphs and bold/italic
        html = md_text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        html = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", html)
        html = re.sub(r"\*(.+?)\*", r"<em>\1</em>", html)
        html = re.sub(r"\n\n+", "</p><p>", html)
        return f"<p>{html}</p>"


def format_table_as_markdown(
    headers: List[str],
    rows: List[List[Any]],
) -> str:
    """
    Format tabular data as a Markdown table string.
    Useful for returning structured data in chat responses.

    Args:
        headers: Column header names.
        rows: List of row data (each row is a list of values).

    Returns:
        Markdown table string.
    """
    if not headers:
        return ""

    # Calculate column widths
    col_widths = [len(str(h)) for h in headers]
    for row in rows:
        for i, cell in enumerate(row):
            if i < len(col_widths):
                col_widths[i] = max(col_widths[i], len(str(cell)))

    # Build header
    header_line = "| " + " | ".join(str(h).ljust(w) for h, w in zip(headers, col_widths)) + " |"
    separator = "| " + " | ".join("-" * w for w in col_widths) + " |"

    # Build rows
    data_lines = []
    for row in rows:
        cells = []
        for i, w in enumerate(col_widths):
            val = str(row[i]) if i < len(row) else ""
            cells.append(val.ljust(w))
        data_lines.append("| " + " | ".join(cells) + " |")

    return "\n".join([header_line, separator] + data_lines)


# ---------------------------------------------------------------------------
# Date / Time Utilities
# ---------------------------------------------------------------------------

def format_datetime(
    dt: Optional[Any] = None,
    fmt: str = "%Y-%m-%d %H:%M:%S",
) -> str:
    """
    Format a datetime to string. If no datetime provided, returns current time.

    Args:
        dt: A datetime object, timestamp (float/int), or None for now.
        fmt: strftime format string.

    Returns:
        Formatted datetime string.
    """
    from datetime import datetime, timezone

    if dt is None:
        dt = datetime.now()
    elif isinstance(dt, (int, float)):
        dt = datetime.fromtimestamp(dt)

    return dt.strftime(fmt)


# ---------------------------------------------------------------------------
# Data Transformation Utilities
# ---------------------------------------------------------------------------

def flatten_dict(d: Dict[str, Any], parent_key: str = "", sep: str = ".") -> Dict[str, Any]:
    """
    Flatten a nested dictionary into a single-level dict with dotted keys.

    Example:
        {"a": {"b": 1, "c": 2}} -> {"a.b": 1, "a.c": 2}
    """
    items: List[tuple] = []
    for k, v in d.items():
        new_key = f"{parent_key}{sep}{k}" if parent_key else k
        if isinstance(v, dict):
            items.extend(flatten_dict(v, new_key, sep).items())
        else:
            items.append((new_key, v))
    return dict(items)


def chunk_list(lst: list, chunk_size: int) -> List[list]:
    """
    Split a list into chunks of a given size.

    Example:
        chunk_list([1,2,3,4,5], 2) -> [[1,2], [3,4], [5]]
    """
    return [lst[i : i + chunk_size] for i in range(0, len(lst), chunk_size)]
