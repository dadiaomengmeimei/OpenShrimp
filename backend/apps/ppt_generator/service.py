"""
Core business logic for PPT Generator.
Handles LLM orchestration, PPT generation, and session management.
"""
from __future__ import annotations

import json
import re
import uuid
from pathlib import Path
from typing import Optional, List, Dict, Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from backend.core.llm_service import chat_completion
from . import config
from .models import Slide, ChatMessage, PPTSession
from .prompts import (
    PPT_GENERATION_SYSTEM,
    PPT_UPDATE_SYSTEM,
    STYLE_DESCRIPTIONS,
)


# In-memory session storage (in production, use a database)
_sessions: Dict[str, PPTSession] = {}


def _create_pptx_file(session: PPTSession) -> Path:
    """Create a PowerPoint file from session data and save it with enhanced styling."""
    filename = f"{session.session_id}.pptx"
    filepath = config.PPT_OUTPUT_DIR.resolve() / filename
    
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Define styles based on session style
    style_config = _get_style_config(session.style)
    
    # Create title slide (first slide)
    if session.slides:
        title_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Add background gradient effect using a shape
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(*style_config["bg_color"])
        bg_shape.line.fill.background()
        
        # Add decorative top bar
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, Inches(0.5)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = RGBColor(*style_config["accent_color"])
        top_bar.line.fill.background()
        
        # Add title with better positioning
        title_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(2.0), Inches(11.333), Inches(2.0)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.add_paragraph()
        title_para.text = session.topic
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.name = 'Microsoft YaHei' if session.language == 'zh' else 'Arial'
        title_para.font.color.rgb = RGBColor(*style_config["title_color"])
        title_para.alignment = PP_ALIGN.CENTER
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Add subtitle with style info
        subtitle_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(3.8), Inches(11.333), Inches(1.0)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.add_paragraph()
        style_name = STYLE_DESCRIPTIONS.get(session.style, session.style)
        lang_name = "中文" if session.language == 'zh' else "English"
        subtitle_para.text = f"{style_name} | {lang_name}"
        subtitle_para.font.size = Pt(20)
        subtitle_para.font.name = 'Microsoft YaHei' if session.language == 'zh' else 'Arial'
        subtitle_para.font.color.rgb = RGBColor(*style_config["subtitle_color"])
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Add decorative bottom bar
        bottom_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(7.0),
            prs.slide_width, Inches(0.5)
        )
        bottom_bar.fill.solid()
        bottom_bar.fill.fore_color.rgb = RGBColor(*style_config["accent_color"])
        bottom_bar.line.fill.background()
    
    # Create content slides
    for idx, slide_data in enumerate(session.slides[1:] if len(session.slides) > 1 else session.slides):
        content_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(content_slide_layout)
        
        # Add background
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(*style_config["bg_color"])
        bg_shape.line.fill.background()
        
        # Add top accent bar
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, Inches(0.4)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = RGBColor(*style_config["accent_color"])
        top_bar.line.fill.background()
        
        # Add slide number indicator
        slide_num_text = f"{idx + 2}" if len(session.slides) > 1 else "1"
        num_box = slide.shapes.add_textbox(
            Inches(12.5), Inches(0.1), Inches(0.8), Inches(0.3)
        )
        num_frame = num_box.text_frame
        num_para = num_frame.add_paragraph()
        num_para.text = slide_num_text
        num_para.font.size = Pt(14)
        num_para.font.color.rgb = RGBColor(*style_config["accent_color"])
        num_para.alignment = PP_ALIGN.RIGHT
        
        # Add slide title with better styling
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.6), Inches(11.733), Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.add_paragraph()
        title_para.text = slide_data.title
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.name = 'Microsoft YaHei' if session.language == 'zh' else 'Arial'
        title_para.font.color.rgb = RGBColor(*style_config["title_color"])
        
        # Add decorative line under title
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(1.7),
            Inches(4.0), Inches(0.1)
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = RGBColor(*style_config["accent_color"])
        line_shape.line.fill.background()
        
        # Add bullet points with better formatting
        if slide_data.content:
            content_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(2.0), Inches(11.733), Inches(4.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for i, point in enumerate(slide_data.content):
                if i == 0:
                    para = content_frame.add_paragraph()
                else:
                    para = content_frame.add_paragraph()
                
                # Add bullet point symbol
                para.text = f"• {point}"
                para.font.size = Pt(22)
                para.font.name = 'Microsoft YaHei' if session.language == 'zh' else 'Arial'
                para.font.color.rgb = RGBColor(*style_config["content_color"])
                para.level = 0
                para.space_before = Pt(10)
                para.space_after = Pt(6)
                para.line_spacing = 1.3
        
        # Add footer with topic
        footer_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(6.8), Inches(11.733), Inches(0.5)
        )
        footer_frame = footer_box.text_frame
        footer_para = footer_frame.add_paragraph()
        footer_para.text = session.topic
        footer_para.font.size = Pt(12)
        footer_para.font.color.rgb = RGBColor(*style_config["subtitle_color"])
        footer_para.alignment = PP_ALIGN.CENTER
        
        # Add speaker notes if available
        if slide_data.notes:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = slide_data.notes
    
    # Save the presentation
    prs.save(str(filepath))
    
    return filepath


def _get_style_config(style: str) -> Dict[str, Any]:
    """Get color scheme and formatting based on style."""
    styles = {
        "professional": {
            "title_color": (25, 25, 112),  # Dark blue
            "subtitle_color": (100, 100, 100),  # Gray
            "content_color": (50, 50, 50),  # Dark gray
            "bg_color": (255, 255, 255),  # White
            "accent_color": (41, 128, 185),  # Blue accent
        },
        "creative": {
            "title_color": (220, 20, 60),  # Crimson
            "subtitle_color": (128, 0, 128),  # Purple
            "content_color": (60, 60, 60),
            "bg_color": (255, 250, 240),  # Floral white
            "accent_color": (230, 126, 34),  # Orange accent
        },
        "minimal": {
            "title_color": (0, 0, 0),  # Black
            "subtitle_color": (128, 128, 128),  # Gray
            "content_color": (64, 64, 64),
            "bg_color": (255, 255, 255),  # White
            "accent_color": (0, 0, 0),  # Black accent
        },
        "academic": {
            "title_color": (0, 100, 0),  # Dark green
            "subtitle_color": (70, 70, 70),
            "content_color": (40, 40, 40),
            "bg_color": (255, 255, 250),  # Snow
            "accent_color": (39, 174, 96),  # Green accent
        },
    }
    return styles.get(style, styles["professional"])


def _parse_llm_response(response: str) -> List[Dict[str, Any]]:
    """
    Parse LLM response to extract JSON array of slides.
    Handles various formats: raw JSON, JSON with markdown, JSON with surrounding text.
    """
    # Remove markdown code blocks if present
    response = re.sub(r'```(?:json)?\s*', '', response)
    response = re.sub(r'```\s*', '', response)
    
    # Try to find JSON array in the response
    json_match = re.search(r'\[.*\]', response, re.DOTALL)
    if json_match:
        try:
            return json.loads(json_match.group())
        except json.JSONDecodeError:
            pass
    
    # Try parsing the whole response as JSON
    try:
        return json.loads(response)
    except json.JSONDecodeError:
        pass
    
    # If all parsing fails, return empty list
    return []


def _validate_and_create_slides(slides_data: List[Dict[str, Any]], min_count: int = 3) -> List[Slide]:
    """Validate slide data and create Slide objects."""
    slides = []
    for slide_data in slides_data:
        if not isinstance(slide_data, dict):
            continue
        slide = Slide(
            title=slide_data.get("title", "Untitled"),
            content=slide_data.get("content", []) if isinstance(slide_data.get("content"), list) else [],
            notes=slide_data.get("notes")
        )
        slides.append(slide)
    
    # Ensure minimum slides
    if len(slides) < min_count:
        while len(slides) < min_count:
            slides.append(Slide(
                title=f"补充内容 {len(slides) + 1}" if min_count == 3 else f"Supplementary Content {len(slides) + 1}",
                content=["可根据需要添加更多内容"],
                notes="此幻灯片可以添加更多内容"
            ))
    
    return slides


async def generate_ppt_from_topic(
    topic: str,
    style: str = "professional",
    language: str = "zh",
    slide_count: Optional[int] = None,
) -> PPTSession:
    """Generate a new PPT presentation from a topic."""
    session_id = uuid.uuid4().hex[:12]
    
    # Build prompt for LLM with better structure requirements
    system_prompt = PPT_GENERATION_SYSTEM.format(language="Chinese" if language == "zh" else "English")
    
    user_prompt = f"""请创建一个关于"{topic}"的 PPT 演示文稿大纲。

要求：
- 风格：{style}
- 语言：{language}
- 幻灯片数量：{slide_count or '自动确定（建议 5-15 页）'}

请返回 JSON 数组格式，每个元素包含：
- title: 幻灯片标题（简洁有力，不超过 10 个字）
- content: 关键点数组（每条不超过 15 个字，3-6 条）
- notes: 可选的演讲者备注

结构要求：
1. 第一页：标题页（主题 + 副标题）
2. 第二页：目录/概览
3. 中间页：主要内容（按逻辑分块）
4. 最后一页：总结/展望

确保内容结构清晰、逻辑连贯。"""
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt},
    ]
    
    # Call LLM
    response = await chat_completion(messages)
    
    # Parse JSON response
    slides_data = _parse_llm_response(response)
    
    if slides_data:
        slides = _validate_and_create_slides(slides_data, config.MIN_SLIDE_COUNT)
    else:
        # Fallback: create a well-structured basic outline
        if language == "zh":
            slides = [
                Slide(title=topic, content=[f"关于{topic}的全面介绍"], notes="标题页"),
                Slide(title="目录", content=["第一部分：背景介绍", "第二部分：核心内容", "第三部分：总结展望"], notes="内容概览"),
                Slide(title="背景介绍", content=["相关概念定义", "发展现状分析", "重要性说明"], notes="介绍背景知识"),
                Slide(title="核心内容", content=["关键点一详细说明", "关键点二详细说明", "关键点三详细说明"], notes="主要内容展示"),
                Slide(title="总结展望", content=["要点回顾", "未来发展方向", "感谢聆听"], notes="总结全文"),
            ]
        else:
            slides = [
                Slide(title=topic, content=[f"Comprehensive introduction to {topic}"], notes="Title slide"),
                Slide(title="Table of Contents", content=["Part 1: Introduction", "Part 2: Main Content", "Part 3: Summary"], notes="Overview"),
                Slide(title="Introduction", content=["Key concepts defined", "Current status analysis", "Importance explained"], notes="Background info"),
                Slide(title="Main Content", content=["Key point one detailed", "Key point two detailed", "Key point three detailed"], notes="Core content"),
                Slide(title="Summary", content=["Key points review", "Future outlook", "Thank you"], notes="Conclusion"),
            ]
    
    # Limit max slides
    if slide_count and len(slides) > slide_count:
        slides = slides[:slide_count]
    elif len(slides) > config.MAX_SLIDE_COUNT:
        slides = slides[:config.MAX_SLIDE_COUNT]
    
    # Generate PPT file FIRST (before creating session)
    temp_session = PPTSession(
        session_id=session_id,
        topic=topic,
        style=style,
        language=language,
        slides=slides,
        chat_history=[],
        ppt_file_path=None
    )
    filepath = _create_pptx_file(temp_session)
    
    # Create session with the file path already set
    session = PPTSession(
        session_id=session_id,
        topic=topic,
        style=style,
        language=language,
        slides=slides,
        chat_history=[
            ChatMessage(role="user", content=f"Generate PPT about {topic}"),
            ChatMessage(role="assistant", content=f"Generated {len(slides)}-page PPT outline")
        ],
        ppt_file_path=str(filepath)
    )
    
    # Store session
    _sessions[session.session_id] = session
    
    return session


async def update_ppt_session(
    session_id: str,
    instruction: str,
    style: Optional[str] = None,
    language: Optional[str] = None,
) -> PPTSession:
    """Update an existing PPT session based on user instruction."""
    if session_id not in _sessions:
        raise ValueError(f"Session {session_id} not found")
    
    session = _sessions[session_id]
    
    # Update style/language if provided
    if style:
        session.style = style
    if language:
        session.language = language
    
    # Build prompt for LLM update
    system_prompt = PPT_UPDATE_SYSTEM.format(
        slide_count=len(session.slides),
        topic=session.topic,
        language="Chinese" if session.language == "zh" else "English"
    )
    
    # Prepare current slides as JSON
    current_slides = [slide.model_dump() for slide in session.slides]
    
    user_prompt = f"""当前用户请求：{instruction}

当前幻灯片内容：
{json.dumps(current_slides, ensure_ascii=False, indent=2)}

请根据用户请求修改幻灯片内容，返回完整的更新后的 JSON 数组。"""
    
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt},
    ]
    
    # Call LLM
    response = await chat_completion(messages)
    
    # Parse JSON response
    slides_data = _parse_llm_response(response)
    
    if slides_data:
        new_slides = _validate_and_create_slides(slides_data, min_count=3)
        session.slides = new_slides
    else:
        # If parsing fails, keep original slides but add a note
        session.chat_history.append(
            ChatMessage(role="user", content=instruction)
        )
        session.chat_history.append(
            ChatMessage(role="assistant", content="抱歉，未能解析更新内容，请重试或提供更明确的指令")
        )
        return session
    
    # Regenerate PPT file with updated content
    filepath = _create_pptx_file(session)
    session.ppt_file_path = str(filepath)
    
    # Update chat history
    session.chat_history.append(
        ChatMessage(role="user", content=instruction)
    )
    session.chat_history.append(
        ChatMessage(role="assistant", content=f"已更新幻灯片，现在包含 {len(session.slides)} 页")
    )
    
    # Store updated session
    _sessions[session_id] = session
    
    return session