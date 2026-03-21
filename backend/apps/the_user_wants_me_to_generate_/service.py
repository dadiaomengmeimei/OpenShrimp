"""Core business logic for PPT Generator app."""

import json
import os
import uuid
import time
import asyncio
from typing import Dict, List, Optional, Any
from datetime import datetime

# Lazy import pptx to avoid module load failure if not installed
_pptx_available = False
Presentation = None
Inches = None
Pt = None
RGBColor = None
PP_ALIGN = None
MSO_ANCHOR = None

def _ensure_pptx():
    """Ensure pptx library is loaded, raise error if not available."""
    global _pptx_available, Presentation, Inches, Pt, RGBColor, PP_ALIGN, MSO_ANCHOR
    if _pptx_available:
        return
    try:
        from pptx import Presentation as _Presentation
        from pptx.util import Inches as _Inches, Pt as _Pt
        from pptx.dml.color import RGBColor as _RGBColor
        from pptx.enum.text import PP_ALIGN as _PP_ALIGN, MSO_ANCHOR as _MSO_ANCHOR
        Presentation = _Presentation
        Inches = _Inches
        Pt = _Pt
        RGBColor = _RGBColor
        PP_ALIGN = _PP_ALIGN
        MSO_ANCHOR = _MSO_ANCHOR
        _pptx_available = True
    except ImportError as e:
        raise ImportError(
            "python-pptx library is required for PPT generation. "
            "Install it with: pip install python-pptx"
        ) from e

from .config import (
    PPT_STORAGE_DIR,
    AVAILABLE_THEMES,
    DEFAULT_THEME,
    LAYOUT_TITLE,
    LAYOUT_TITLE_AND_CONTENT,
    LAYOUT_BLANK,
    SESSION_TIMEOUT,
)
from .models import PPTOutline, SlideContent, SlideType, PPTSession
from .prompts import SYSTEM_PROMPT_OUTLINE, SYSTEM_PROMPT_MODIFY, SYSTEM_PROMPT_OUTLINE_WITH_FILE
from .file_parser import ParsedFile, combine_parsed_files

# In-memory session storage (in production, use Redis or database)
_sessions: Dict[str, PPTSession] = {}


def _get_llm_service():
    """Lazy import LLM service to avoid circular imports."""
    from backend.core.llm_service import chat_completion
    return chat_completion


def ensure_storage_dir():
    """Ensure storage directory exists."""
    if not os.path.exists(PPT_STORAGE_DIR):
        os.makedirs(PPT_STORAGE_DIR, exist_ok=True)


def create_session() -> PPTSession:
    """Create a new PPT generation session."""
    session_id = str(uuid.uuid4())
    now = time.time()
    session = PPTSession(
        session_id=session_id,
        created_at=now,
        last_updated=now,
        theme=DEFAULT_THEME,
    )
    _sessions[session_id] = session
    return session


def get_session(session_id: str) -> Optional[PPTSession]:
    """Get session by ID."""
    session = _sessions.get(session_id)
    if session and time.time() - session.last_updated > SESSION_TIMEOUT:
        del _sessions[session_id]
        return None
    return session


def update_session(session: PPTSession):
    """Update session timestamp and save."""
    session.last_updated = time.time()
    _sessions[session.session_id] = session


def cleanup_old_sessions():
    """Remove expired sessions."""
    now = time.time()
    expired = [
        sid for sid, s in _sessions.items()
        if now - s.last_updated > SESSION_TIMEOUT
    ]
    for sid in expired:
        del _sessions[sid]


def parse_json_from_llm(text: str) -> Optional[Dict]:
    """Extract and parse JSON from LLM response."""
    # Try to find JSON block
    start_idx = text.find('{')
    end_idx = text.rfind('}')
    
    if start_idx == -1 or end_idx == -1 or end_idx <= start_idx:
        return None
    
    json_str = text[start_idx:end_idx + 1]
    
    try:
        return json.loads(json_str)
    except json.JSONDecodeError:
        # Try cleaning up common issues
        try:
            # Replace single quotes with double quotes
            cleaned = json_str.replace("'", '"')
            return json.loads(cleaned)
        except json.JSONDecodeError:
            return None


def _create_fallback_outline(topic: str, theme: str, slide_count: int) -> Dict:
    """Create a fallback outline when LLM fails or times out."""
    slides = [
        {"type": "title", "title": topic, "subtitle": "精彩演示文稿"},
        {"type": "content", "title": "概述", "content": f"本演示文稿介绍{topic}的核心内容与价值"},
    ]
    
    # Add content slides
    content_slides = min(slide_count - 3, 5)  # Leave room for title and end
    for i in range(content_slides):
        slides.append({
            "type": "content",
            "title": f"第{i+1}部分",
            "content": f"• {topic}的重要方面{i+1}\n• 关键要点与洞察\n• 实践应用与案例"
        })
    
    # Add end slide
    slides.append({"type": "end", "title": "谢谢观看", "content": "欢迎交流讨论"})
    
    return {
        "title": topic,
        "subtitle": "精彩演示文稿",
        "theme": theme,
        "total_slides": len(slides),
        "slides": slides,
    }


async def generate_outline(topic: str, theme: str = DEFAULT_THEME, slide_count: int = 8) -> PPTOutline:
    """Generate PPT outline using LLM."""
    chat_completion = _get_llm_service()
    
    messages = [
        {"role": "system", "content": SYSTEM_PROMPT_OUTLINE},
        {"role": "user", "content": f"请为以下主题生成PPT大纲，使用{theme}主题，共{slide_count}页：\n\n{topic}"}
    ]
    
    return await _generate_outline_with_messages(messages, topic, theme, slide_count)


async def generate_outline_from_content(content: str, theme: str = DEFAULT_THEME, slide_count: int = 8, user_request: str = "") -> PPTOutline:
    """Generate PPT outline from file content using LLM.
    
    Args:
        content: The combined content from uploaded files and/or user text
        theme: Theme to use for the presentation
        slide_count: Number of slides to generate
        user_request: Original user request text (for fallback title)
    """
    chat_completion = _get_llm_service()
    
    # Build the prompt with file content
    prompt_content = f"""请根据以下内容生成PPT大纲，使用{theme}主题，共{slide_count}页。

用户原始请求：{user_request if user_request else "根据文件内容生成PPT"}

=== 参考内容开始 ===
{content}
=== 参考内容结束 ===

请仔细阅读上述内容，提取关键信息，生成一份结构清晰、内容专业的PPT大纲。"""
    
    messages = [
        {"role": "system", "content": SYSTEM_PROMPT_OUTLINE_WITH_FILE},
        {"role": "user", "content": prompt_content}
    ]
    
    try:
        # Use timeout and limit tokens to ensure faster response
        response = await asyncio.wait_for(
            chat_completion(messages, max_tokens=2000),
            timeout=60.0  # 60 second timeout for outline generation
        )
    except asyncio.TimeoutError:
        # Use fallback outline on timeout
        data = _create_fallback_outline(user_request or "演示文稿", theme, slide_count)
    except Exception as e:
        # Use fallback outline on any error
        data = _create_fallback_outline(user_request or "演示文稿", theme, slide_count)
    else:
        # Parse LLM response
        data = parse_json_from_llm(response)
        if not data:
            # Fallback: create a simple outline
            data = _create_fallback_outline(user_request or "演示文稿", theme, slide_count)
    
    # Parse slides
    slides = []
    for slide_data in data.get("slides", []):
        try:
            slide_type = SlideType(slide_data.get("type", "content"))
        except ValueError:
            slide_type = SlideType.CONTENT
        
        slides.append(SlideContent(
            type=slide_type,
            title=slide_data.get("title", "Untitled"),
            content=slide_data.get("content"),
            subtitle=slide_data.get("subtitle"),
            notes=slide_data.get("notes"),
        ))
    
    return PPTOutline(
        title=data.get("title", user_request or "演示文稿"),
        subtitle=data.get("subtitle"),
        theme=data.get("theme", theme),
        total_slides=data.get("total_slides", len(slides)),
        slides=slides,
    )


async def modify_outline(session: PPTSession, instruction: str) -> PPTOutline:
    """Modify existing outline based on user instruction."""
    if not session.outline:
        raise ValueError("No existing outline to modify")
    
    chat_completion = _get_llm_service()
    
    # Build outline summary for context
    outline_summary = []
    for i, slide in enumerate(session.outline.slides, 1):
        outline_summary.append(f"{i}. [{slide.type.value}] {slide.title}")
    
    prompt = SYSTEM_PROMPT_MODIFY.format(
        title=session.outline.title,
        theme=session.outline.theme,
        slide_count=len(session.outline.slides),
        outline="\n".join(outline_summary),
        instruction=instruction,
    )
    
    messages = [
        {"role": "system", "content": prompt},
        {"role": "user", "content": instruction}
    ]
    
    try:
        response = await asyncio.wait_for(
            chat_completion(messages, max_tokens=2000),
            timeout=60.0
        )
    except asyncio.TimeoutError:
        # On timeout, try simple modifications or return original
        return _apply_simple_modification(session.outline, instruction)
    except Exception:
        return _apply_simple_modification(session.outline, instruction)
    
    data = parse_json_from_llm(response)
    if not data:
        # If parsing fails, try simple modifications
        return _apply_simple_modification(session.outline, instruction)
    
    # Parse modified slides
    slides = []
    for slide_data in data.get("slides", []):
        try:
            slide_type = SlideType(slide_data.get("type", "content"))
        except ValueError:
            slide_type = SlideType.CONTENT
        
        slides.append(SlideContent(
            type=slide_type,
            title=slide_data.get("title", "Untitled"),
            content=slide_data.get("content"),
            subtitle=slide_data.get("subtitle"),
            notes=slide_data.get("notes"),
        ))
    
    return PPTOutline(
        title=data.get("title", session.outline.title),
        subtitle=data.get("subtitle", session.outline.subtitle),
        theme=data.get("theme", session.outline.theme),
        total_slides=data.get("total_slides", len(slides)),
        slides=slides,
    )


def _apply_simple_modification(outline: PPTOutline, instruction: str) -> PPTOutline:
    """Apply simple modifications without LLM."""
    instruction_lower = instruction.lower()
    
    # Handle simple style changes
    for theme_key in AVAILABLE_THEMES.keys():
        if theme_key in instruction_lower or AVAILABLE_THEMES[theme_key]["name"] in instruction:
            outline.theme = theme_key
            return outline
    
    # Return original if we can't parse
    return outline


def hex_to_rgb(hex_color: str) -> tuple:
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def apply_theme_to_slide(slide, theme_config: Dict, is_title_slide: bool = False):
    """Apply theme colors to slide."""
    _ensure_pptx()
    primary_color = theme_config.get("primary_color", "1F4E79")
    secondary_color = theme_config.get("secondary_color", "2E75B6")
    
    # Apply to all shapes with text
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Set font color
                    r, g, b = hex_to_rgb(primary_color)
                    run.font.color.rgb = RGBColor(r, g, b)


def create_ppt_file(outline: PPTOutline, session_id: str) -> str:
    """Create PPT file from outline."""
    _ensure_pptx()  # Ensure pptx is available before using
    
    ensure_storage_dir()
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    theme_config = AVAILABLE_THEMES.get(outline.theme, AVAILABLE_THEMES[DEFAULT_THEME])
    
    for i, slide_content in enumerate(outline.slides):
        # Select layout based on type
        if slide_content.type == SlideType.TITLE:
            layout = prs.slide_layouts[LAYOUT_TITLE]
        elif slide_content.type == SlideType.SECTION:
            layout = prs.slide_layouts[LAYOUT_BLANK]
        else:
            layout = prs.slide_layouts[LAYOUT_TITLE_AND_CONTENT]
        
        slide = prs.slides.add_slide(layout)
        
        # Apply theme background for dark theme
        if outline.theme == "dark":
            background = slide.background
            fill = background.fill
            fill.solid()
            bg_color = theme_config.get("bg_color", "1a1a2e")
            r, g, b = hex_to_rgb(bg_color)
            fill.fore_color.rgb = RGBColor(r, g, b)
        
        # Set title
        if slide.shapes.title:
            title = slide.shapes.title
            title.text = slide_content.title
            
            # Style title
            title_frame = title.text_frame
            title_frame.paragraphs[0].font.size = Pt(44 if slide_content.type == SlideType.TITLE else 32)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.name = theme_config.get("font_title", "微软雅黑")
            
            # Set text color based on theme
            if outline.theme == "dark":
                for run in title_frame.paragraphs[0].runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)
            else:
                r, g, b = hex_to_rgb(theme_config.get("primary_color", "1F4E79"))
                for run in title_frame.paragraphs[0].runs:
                    run.font.color.rgb = RGBColor(r, g, b)
        
        # Set subtitle for title slide
        if slide_content.type == SlideType.TITLE and slide_content.subtitle:
            # Find subtitle placeholder
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 2:  # SUBTITLE
                    shape.text = slide_content.subtitle
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.size = Pt(24)
                        paragraph.font.name = theme_config.get("font_body", "微软雅黑")
                        if outline.theme == "dark":
                            for run in paragraph.runs:
                                run.font.color.rgb = RGBColor(200, 200, 200)
                    break
        
        # Set content
        if slide_content.content and len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = slide_content.content
            
            # Style content
            text_frame = content_placeholder.text_frame
            text_frame.word_wrap = True
            
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.name = theme_config.get("font_body", "微软雅黑")
                paragraph.space_after = Pt(12)
                
                # Set text color
                if outline.theme == "dark":
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(230, 230, 230)
                else:
                    r, g, b = hex_to_rgb(theme_config.get("secondary_color", "404040"))
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(r, g, b)
                
                # Handle bullet points
                if slide_content.type == SlideType.BULLETS:
                    paragraph.level = 0
    
    # Save file
    filename = f"ppt_{session_id}_{int(time.time())}.pptx"
    filepath = os.path.join(PPT_STORAGE_DIR, filename)
    prs.save(filepath)
    
    return filepath


def format_outline_for_display(outline: PPTOutline) -> str:
    """Format outline for display in chat."""
    lines = [
        f"**{outline.title}**",
        "",
        "📑 **大纲预览**：",
    ]
    
    for i, slide in enumerate(outline.slides, 1):
        icon = "📌" if slide.type == SlideType.TITLE else \
               "📋" if slide.type == SlideType.SECTION else \
               "📝" if slide.type == SlideType.CONTENT else \
               "🔚" if slide.type == SlideType.END else "•"
        lines.append(f"{icon} 第{i}页：{slide.title}")
    
    lines.append("")
    lines.append(f"🎨 **主题**：{AVAILABLE_THEMES.get(outline.theme, {}).get('name', outline.theme)}")
    lines.append(f"📊 **页数**：{len(outline.slides)}页")
    
    return "\n".join(lines)


def get_theme_list() -> List[Dict]:
    """Get list of available themes for API."""
    return [
        {"key": key, **config}
        for key, config in AVAILABLE_THEMES.items()
    ]