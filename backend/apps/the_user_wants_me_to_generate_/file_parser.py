"""File parsing utilities for PPT Generator app.

Supports extracting text content from various file types:
- Text files (.txt, .md, .csv)
- Word documents (.docx)
- PDF files (.pdf)
- Excel files (.xlsx, .xls)
- PowerPoint files (.pptx)
"""

import os
from typing import Optional, Dict, Any, List
from dataclasses import dataclass


@dataclass
class ParsedFile:
    """Represents a parsed file with its metadata and content."""
    filename: str
    file_type: str
    content: str
    metadata: Dict[str, Any]
    summary: str = ""


# Supported file types and their MIME types/extensions
SUPPORTED_EXTENSIONS = {
    '.txt': 'text',
    '.md': 'text',
    '.csv': 'csv',
    '.docx': 'docx',
    '.pdf': 'pdf',
    '.xlsx': 'excel',
    '.xls': 'excel',
    '.pptx': 'pptx',
}


def detect_file_type(file_path: str) -> Optional[str]:
    """Detect file type based on extension."""
    ext = os.path.splitext(file_path)[1].lower()
    return SUPPORTED_EXTENSIONS.get(ext)


def parse_text_file(file_path: str) -> str:
    """Parse plain text file."""
    encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    # Fallback: read as binary and decode with errors='ignore'
    with open(file_path, 'rb') as f:
        return f.read().decode('utf-8', errors='ignore')


def parse_docx_file(file_path: str) -> str:
    """Parse Word document using file_toolkit."""
    try:
        from backend.core.file_toolkit import parse_docx
        return parse_docx(file_path)
    except Exception as e:
        return f"[Error parsing DOCX: {str(e)}]"


def parse_pdf_file(file_path: str) -> str:
    """Parse PDF file using file_toolkit."""
    try:
        from backend.core.file_toolkit import parse_pdf
        return parse_pdf(file_path)
    except Exception as e:
        return f"[Error parsing PDF: {str(e)}]"


def parse_excel_file(file_path: str) -> str:
    """Parse Excel file and convert to formatted text.
    
    Returns a formatted string with all sheets' data.
    """
    try:
        from backend.core.file_toolkit import parse_excel
        
        result = parse_excel(file_path)
        if not result:
            return "[Empty Excel file or unable to parse]"
        
        # Handle both single sheet and multiple sheets
        if isinstance(result, dict) and 'rows' in result:
            # Single sheet result
            return _format_excel_data(result)
        elif isinstance(result, dict):
            # Multiple sheets - format each
            sections = []
            for sheet_name, sheet_data in result.items():
                sections.append(f"\n--- Sheet: {sheet_name} ---\n")
                sections.append(_format_excel_data(sheet_data))
            return "\n".join(sections)
        else:
            return str(result)
    except Exception as e:
        return f"[Error parsing Excel: {str(e)}]"


def _format_excel_data(sheet_data: Dict[str, Any]) -> str:
    """Format Excel sheet data as readable text."""
    lines = []
    
    # Add headers
    headers = sheet_data.get('headers', [])
    if headers:
        lines.append(" | ".join(str(h) for h in headers))
        lines.append("-" * (len(lines[0]) if lines else 20))
    
    # Add rows
    rows = sheet_data.get('rows', [])
    for row in rows[:50]:  # Limit to first 50 rows
        if isinstance(row, dict):
            row_values = [str(row.get(h, '')) for h in headers]
        else:
            row_values = [str(cell) for cell in row]
        lines.append(" | ".join(row_values))
    
    if len(rows) > 50:
        lines.append(f"\n... ({len(rows) - 50} more rows)")
    
    return "\n".join(lines)


def parse_csv_file(file_path: str) -> str:
    """Parse CSV file using file_toolkit."""
    try:
        import csv
        lines = []
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                if i >= 50:  # Limit rows
                    lines.append(f"\n... (more rows)")
                    break
                lines.append(" | ".join(row))
        return "\n".join(lines)
    except Exception as e:
        return f"[Error parsing CSV: {str(e)}]"


def parse_pptx_file(file_path: str) -> str:
    """Extract text from PowerPoint file."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = [f"\n--- Slide {i} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            slides_text.append("\n".join(slide_text))
        
        return "\n".join(slides_text)
    except Exception as e:
        return f"[Error parsing PPTX: {str(e)}]"


def parse_file(file_path: str) -> ParsedFile:
    """Parse a file and return its content.
    
    Args:
        file_path: Path to the file
        
    Returns:
        ParsedFile object with extracted content
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    
    filename = os.path.basename(file_path)
    file_type = detect_file_type(file_path)
    
    if not file_type:
        # Try to read as text anyway
        content = parse_text_file(file_path)
        file_type = 'unknown'
    elif file_type == 'text':
        content = parse_text_file(file_path)
    elif file_type == 'docx':
        content = parse_docx_file(file_path)
    elif file_type == 'pdf':
        content = parse_pdf_file(file_path)
    elif file_type == 'excel':
        content = parse_excel_file(file_path)
    elif file_type == 'csv':
        content = parse_csv_file(file_path)
    elif file_type == 'pptx':
        content = parse_pptx_file(file_path)
    else:
        content = parse_text_file(file_path)
    
    # Generate summary
    summary = _generate_summary(content, file_type)
    
    # Extract metadata
    metadata = {
        'file_size': os.path.getsize(file_path),
        'line_count': content.count('\n') + 1,
        'char_count': len(content),
    }
    
    return ParsedFile(
        filename=filename,
        file_type=file_type,
        content=content,
        metadata=metadata,
        summary=summary
    )


def _generate_summary(content: str, file_type: str) -> str:
    """Generate a brief summary of the file content."""
    lines = content.strip().split('\n')
    non_empty_lines = [l for l in lines if l.strip()]
    
    # First few non-empty lines as preview
    preview_lines = non_empty_lines[:5]
    preview = '\n'.join(preview_lines)
    
    if len(preview) > 300:
        preview = preview[:300] + "..."
    
    return f"[{file_type.upper()} file, {len(non_empty_lines)} lines]\nPreview:\n{preview}"


def extract_files_from_messages(messages: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Extract file attachments from chat messages.
    
    The platform may include file information in messages in various formats:
    - file_path: direct file path
    - files: list of file objects
    - attachments: list of attachments
    
    Returns:
        List of file info dicts with 'path' and optional 'name' keys
    """
    files = []
    
    for msg in messages:
        if msg.get('role') != 'user':
            continue
        
        content = msg.get('content', '')
        
        # Check for file_path in message
        if 'file_path' in msg:
            files.append({
                'path': msg['file_path'],
                'name': msg.get('file_name', os.path.basename(msg['file_path']))
            })
        
        # Check for files list
        if 'files' in msg and isinstance(msg['files'], list):
            for f in msg['files']:
                if isinstance(f, dict) and 'path' in f:
                    files.append({
                        'path': f['path'],
                        'name': f.get('name', os.path.basename(f['path']))
                    })
                elif isinstance(f, str):
                    files.append({'path': f, 'name': os.path.basename(f)})
        
        # Check for attachments list
        if 'attachments' in msg and isinstance(msg['attachments'], list):
            for att in msg['attachments']:
                if isinstance(att, dict) and 'path' in att:
                    files.append({
                        'path': att['path'],
                        'name': att.get('name', os.path.basename(att['path']))
                    })
        
        # Try to detect file paths in content (for backward compatibility)
        # Look for common file patterns
        import re
        file_patterns = [
            r'/[^\s]*\.(?:txt|md|csv|docx|pdf|xlsx?|pptx?)',
            r'uploads?/[^\s]*\.(?:txt|md|csv|docx|pdf|xlsx?|pptx?)',
        ]
        for pattern in file_patterns:
            matches = re.findall(pattern, content, re.IGNORECASE)
            for match in matches:
                if os.path.exists(match):
                    files.append({'path': match, 'name': os.path.basename(match)})
    
    # Remove duplicates based on path
    seen = set()
    unique_files = []
    for f in files:
        if f['path'] not in seen:
            seen.add(f['path'])
            unique_files.append(f)
    
    return unique_files


def combine_parsed_files(parsed_files: List[ParsedFile]) -> str:
    """Combine multiple parsed files into a single context string.
    
    Args:
        parsed_files: List of ParsedFile objects
        
    Returns:
        Combined content string with file separators
    """
    if not parsed_files:
        return ""
    
    if len(parsed_files) == 1:
        return parsed_files[0].content
    
    sections = []
    for pf in parsed_files:
        sections.append(f"\n{'='*60}")
        sections.append(f"FILE: {pf.filename}")
        sections.append(f"TYPE: {pf.file_type}")
        sections.append('='*60)
        sections.append(pf.content)
    
    return "\n".join(sections)


def format_file_summary(parsed_files: List[ParsedFile]) -> str:
    """Format a summary of parsed files for display to user.
    
    Args:
        parsed_files: List of ParsedFile objects
        
    Returns:
        Formatted summary string
    """
    if not parsed_files:
        return ""
    
    lines = ["📎 **已上传文件：**"]
    for pf in parsed_files:
        size_kb = pf.metadata['file_size'] / 1024
        lines.append(f"  • {pf.filename} ({pf.file_type}, {size_kb:.1f}KB, {pf.metadata['line_count']}行)")
    
    return "\n".join(lines)